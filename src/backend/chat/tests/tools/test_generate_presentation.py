"""
Tests for generate_presentation.

Real components: Django ORM (factory-built conversation), real default_storage,
the real pptx template and builder, real RunContext + ContextDeps.

The only thing mocked is the PresentationAgent's LLM (via FunctionModel) - the
standard pydantic-ai idiom for driving deterministic model output.
"""

import re
from io import BytesIO
from unittest import mock
from urllib.parse import parse_qs, urlparse

from django.core.files.storage import default_storage

import pptx
import pytest
from asgiref.sync import sync_to_async
from pydantic_ai import ModelResponse, RunContext
from pydantic_ai.exceptions import ModelRetry
from pydantic_ai.messages import ToolCallPart
from pydantic_ai.models.function import FunctionModel
from pydantic_ai.usage import RunUsage

from chat.agents.presentation import PresentationAgent
from chat.clients.schema import ContextDeps
from chat.factories import ChatConversationFactory, UserFactory
from chat.llm_configuration import LLModel, LLMProvider
from chat.tools.generate_presentation import (
    DOWNLOAD_URL_EXPIRATION,
    build_object_name,
    generate_presentation,
)

# transaction=True is required so writes done via sync_to_async (which run on
# threadpool connections distinct from the test's wrapping transaction) commit
# and are flushed via TRUNCATE between tests instead of leaking across them.
pytestmark = pytest.mark.django_db(transaction=True)


@pytest.fixture(autouse=True)
def fixture_presentation_agent_config(settings):
    """Configure the LLM model used by PresentationAgent."""
    settings.LLM_CONFIGURATIONS = {
        settings.LLM_DEFAULT_MODEL_HRID: LLModel(
            hrid="mistral-model",
            model_name="mistral-7b-instruct-v0.1",
            human_readable_name="Mistral 7B Instruct",
            profile=None,
            provider=LLMProvider(
                hrid="mistral",
                kind="mistral",
                base_url="https://api.mistral.ai/v1",
                api_key="testkey",
            ),
            is_active=True,
            system_prompt="direct",
            tools=[],
        ),
    }


DECK = {
    "title": "Sobriété énergétique",
    "slides": [
        {"type": "cover", "title": "Sobriété énergétique"},
        {
            "type": "title_one_column",
            "title": "Constats",
            "content": "- Premier **constat**\n- Second constat",
            "slide_notes": "Insister sur le premier point",
        },
    ],
}


def model_returns_deck(deck):
    """Build a FunctionModel callback answering with `deck` as structured output."""

    def respond(_messages, info):
        return ModelResponse(parts=[ToolCallPart(tool_name=info.output_tools[0].name, args=deck)])

    return respond


@sync_to_async
def setup_context(max_retries=2):
    """Build a conversation and a real RunContext in one sync block."""
    user = UserFactory()
    conversation = ChatConversationFactory(owner=user)
    ctx = RunContext(
        model="test",
        usage=RunUsage(input_tokens=0, output_tokens=0),
        deps=ContextDeps(conversation=conversation, user=user),
        max_retries=max_retries,
        retries={},
        tool_name="generate_presentation",
    )
    return ctx, conversation, user


async def run_tool(ctx, brief="Un support sur la sobriété énergétique", deck=None):
    """Run the tool with the presentation agent's LLM stubbed out."""
    agent = PresentationAgent()
    with agent.override(model=FunctionModel(model_returns_deck(deck or DECK))):
        with mock.patch("chat.tools.generate_presentation.PresentationAgent", return_value=agent):
            return await generate_presentation(ctx, brief=brief)


def object_key(url):
    """Recover the S3 object key from a presigned URL."""
    path = urlparse(url).path
    return path.split(f"/{default_storage.bucket_name}/", 1)[1]


@pytest.mark.asyncio
async def test_returns_a_signed_time_limited_link():
    """The deck is handed back as a signed URL scoped to the conversation."""
    ctx, conversation, _user = await setup_context()

    result = await run_tool(ctx)

    url = result.metadata["url"]
    assert url in result.return_value

    key = object_key(url)
    assert key.startswith(f"{conversation.pk}/attachments/")
    # Readable slug from the deck title, plus a random hex, not a bare UUID.
    assert re.fullmatch(r"sobriete_energetique_[0-9a-f]{8}\.pptx", key.rsplit("/", 1)[1])

    query = parse_qs(urlparse(url).query)
    assert query["X-Amz-Signature"]
    assert query["X-Amz-Expires"] == [str(DOWNLOAD_URL_EXPIRATION)]


@pytest.mark.asyncio
async def test_the_stored_file_is_a_readable_deck():
    """What lands in storage opens as a deck holding the generated slides."""
    ctx, _conversation, _user = await setup_context()

    result = await run_tool(ctx)

    key = object_key(result.metadata["url"])
    blob = await sync_to_async(lambda: default_storage.open(key).read())()
    deck = pptx.Presentation(BytesIO(blob))

    assert len(deck.slides) == 2
    assert deck.slides[0].slide_layout.name == "Titre et sous-titre"
    assert deck.slides[1].notes_slide.notes_text_frame.text == "Insister sur le premier point"


@pytest.mark.asyncio
async def test_no_attachment_row_is_created():
    """Generated decks live in storage only; nothing is recorded in database."""
    ctx, conversation, _user = await setup_context()

    await run_tool(ctx)

    assert await conversation.attachments.acount() == 0


@pytest.mark.asyncio
async def test_tells_the_model_not_to_restate_the_deck():
    """The return value steers the model away from repeating the slides in chat."""
    ctx, _conversation, _user = await setup_context()

    result = await run_tool(ctx)

    assert "2 slides" in result.return_value
    assert "do not restate" in result.return_value.lower()


@pytest.mark.asyncio
async def test_an_empty_brief_is_sent_back_to_the_model():
    """An empty brief is a model mistake it can correct while retries remain."""
    ctx, _conversation, _user = await setup_context(max_retries=2)

    with pytest.raises(ModelRetry, match="brief is empty"):
        await run_tool(ctx, brief="   ")


@pytest.mark.asyncio
async def test_an_empty_brief_soft_fails_once_retries_run_out():
    """
    Out of retries, the tool returns guidance instead of raising.

    The tool is registered with `retries=1`, so this is what the model actually
    sees on a second empty brief: a message to relay, rather than an exception
    that would let it answer from its own knowledge.
    """
    ctx, _conversation, _user = await setup_context(max_retries=1)

    result = await run_tool(ctx, brief="   ")

    assert "brief is empty" in result


@pytest.mark.asyncio
async def test_nothing_is_written_when_the_brief_is_empty():
    """A rejected brief writes no orphan file to storage."""
    ctx, conversation, _user = await setup_context(max_retries=2)

    with pytest.raises(ModelRetry):
        await run_tool(ctx, brief="")

    _dirs, files = await sync_to_async(default_storage.listdir)(f"{conversation.pk}/attachments")
    assert files == []


@pytest.mark.parametrize(
    "title,expected_slug",
    [
        ("Sobriété énergétique", "sobriete_energetique"),
        ("Bilan 2026 / T1 : résultats !", "bilan_2026_t1_resultats"),
        ("   ", "presentation"),  # empty title falls back
        ("///", "presentation"),  # no alphanumerics falls back
        ("A" * 120, "a" * 60),  # long titles are capped
    ],
)
def test_build_object_name_is_readable_and_unguessable(title, expected_slug):
    """The file name is a snake_case slug plus a random hex, ending in .pptx."""
    name = build_object_name(title)
    assert re.fullmatch(rf"{expected_slug}_[0-9a-f]{{8}}\.pptx", name)


def test_build_object_name_hex_suffix_varies():
    """Two calls for the same title yield different, non-guessable names."""
    first = build_object_name("Deck")
    second = build_object_name("Deck")
    assert first != second
