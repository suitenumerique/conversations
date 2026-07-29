"""
Tests for the presentation builder.

Real components: the bundled pptx template and python-pptx. Nothing is mocked;
decks are built in memory and read back with python-pptx to assert what a user
would actually open.
"""

import dataclasses
from io import BytesIO

import pptx
import pytest

from chat.file_generation.builder import (
    PresentationBuildError,
    build_presentation,
    check_template,
)
from chat.file_generation.entities import Presentation, Slide, SlideType
from chat.file_generation.templates import GENERIC_TEMPLATE


def build_deck(slides):
    """Build a deck from slides and reopen it as python-pptx would."""
    blob = build_presentation(GENERIC_TEMPLATE, Presentation(title="Deck", slides=slides))
    return pptx.Presentation(BytesIO(blob))


def shape_texts(slide):
    """Return the non-empty texts carried by a slide, in shape order."""
    return [
        shape.text_frame.text
        for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text.strip()
    ]


def test_check_template_accepts_the_bundled_template():
    """The declared layouts and placeholder indices match the shipped file."""
    check_template(GENERIC_TEMPLATE)


def test_every_slide_type_has_a_layout():
    """Every slide type the model may emit is bound to a layout."""
    for slide_type in SlideType:
        assert GENERIC_TEMPLATE.get_layout(slide_type) is not None


def test_check_template_rejects_an_unknown_layout():
    """A layout name absent from the pptx file is reported, not silently skipped."""
    broken_layout = dataclasses.replace(GENERIC_TEMPLATE.layouts[0], layout_name="Nonexistent")
    broken = dataclasses.replace(GENERIC_TEMPLATE, layouts=(broken_layout,))

    with pytest.raises(PresentationBuildError, match="Nonexistent"):
        check_template(broken)


def test_check_template_rejects_an_out_of_range_placeholder():
    """A placeholder index beyond the slide's shapes is reported."""
    broken_layout = dataclasses.replace(
        GENERIC_TEMPLATE.layouts[0], shape_index_to_field={99: "title"}
    )
    broken = dataclasses.replace(GENERIC_TEMPLATE, layouts=(broken_layout,))

    with pytest.raises(PresentationBuildError, match="placeholder index 99"):
        check_template(broken)


@pytest.mark.parametrize(
    "slide,expected_layout,expected_texts",
    [
        (
            Slide(type=SlideType.COVER, title="Titre"),
            "Titre et sous-titre",
            {"Titre"},
        ),
        (
            Slide(type=SlideType.SECTION, title="Partie", subtitle="Sous-titre"),
            "Chapitre",
            {"Partie", "Sous-titre"},
        ),
        (
            Slide(type=SlideType.TITLE_ONE_COLUMN, title="Titre", content="Corps"),
            "Titre et textes 1 colonne",
            {"Titre", "Corps"},
        ),
        (
            Slide(type=SlideType.TITLE_TWO_COLUMNS, title="Titre", left="G", right="D"),
            "Titre et textes 2 colonnes",
            {"Titre", "G", "D"},
        ),
        (
            Slide(
                type=SlideType.TITLE_THREE_COLUMNS,
                title="Titre",
                left="A",
                center="B",
                right="C",
            ),
            "Titre et textes 3 colonnes",
            {"Titre", "A", "B", "C"},
        ),
    ],
)
def test_each_slide_type_uses_its_layout(slide, expected_layout, expected_texts):
    """Every slide type lands on its bound layout with its fields filled."""
    deck = build_deck([slide])

    assert len(deck.slides) == 1
    assert deck.slides[0].slide_layout.name == expected_layout
    assert set(shape_texts(deck.slides[0])) == expected_texts


def test_slides_keep_their_order():
    """Slides appear in the order the spec declares them."""
    deck = build_deck(
        [
            Slide(type=SlideType.COVER, title="Un"),
            Slide(type=SlideType.SECTION, title="Deux"),
            Slide(type=SlideType.COVER, title="Trois"),
        ]
    )

    assert [shape_texts(slide)[0] for slide in deck.slides] == ["Un", "Deux", "Trois"]


def test_presenter_notes_are_attached():
    """Slide notes land on the slide's notes page."""
    deck = build_deck([Slide(type=SlideType.COVER, title="Titre", slide_notes="À dire à l'oral")])

    assert deck.slides[0].notes_slide.notes_text_frame.text == "À dire à l'oral"


def test_slides_without_notes_are_left_alone():
    """No notes page is created when a slide carries no notes."""
    deck = build_deck([Slide(type=SlideType.COVER, title="Titre")])

    assert not deck.slides[0].has_notes_slide


def test_unbound_slide_type_is_reported():
    """A type with no layout binding fails loudly instead of silently."""
    cover_only = dataclasses.replace(GENERIC_TEMPLATE, layouts=(GENERIC_TEMPLATE.layouts[0],))
    spec = Presentation(title="Deck", slides=[Slide(type=SlideType.SECTION, title="x")])

    with pytest.raises(PresentationBuildError, match="No layout is bound"):
        build_presentation(cover_only, spec)


def test_missing_template_file_is_reported():
    """A template path that cannot be opened raises a build error."""
    missing = dataclasses.replace(GENERIC_TEMPLATE, path="/nonexistent/template.pptx")
    spec = Presentation(title="x", slides=[Slide(type=SlideType.COVER)])

    with pytest.raises(PresentationBuildError, match="Could not open"):
        build_presentation(missing, spec)


def test_a_deck_needs_at_least_one_slide():
    """An empty deck is rejected at validation time."""
    with pytest.raises(ValueError):
        Presentation(title="Deck", slides=[])
