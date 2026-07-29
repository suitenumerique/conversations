"""
Tests for the Markdown-to-pptx renderer.

Real components: the bundled pptx template, python-markdown and python-pptx.
Content is rendered into a real placeholder taken from the template so the
assertions run against the same shape types production uses.
"""

import pptx
import pytest

from chat.file_generation.renderer import render
from chat.file_generation.templates.generic import TEMPLATE_PATH

BODY_LAYOUT = "Titre et textes 1 colonne"
BODY_PLACEHOLDER_INDEX = 1


@pytest.fixture(name="slide")
def slide_fixture():
    """A real slide built from the bundled template's body layout."""
    deck = pptx.Presentation(str(TEMPLATE_PATH))
    return deck.slides.add_slide(deck.slide_layouts.get_by_name(BODY_LAYOUT))


@pytest.fixture(name="placeholder")
def placeholder_fixture(slide):
    """The body placeholder that renders receive."""
    return slide.shapes[BODY_PLACEHOLDER_INDEX]


def table_of(slide):
    """Return the single table shape rendered onto the slide."""
    return next(shape for shape in slide.shapes if shape.has_table).table


def runs_of(placeholder):
    """Flatten every run of the placeholder across its paragraphs."""
    return [run for paragraph in placeholder.text_frame.paragraphs for run in paragraph.runs]


def test_plain_text_is_rendered(placeholder):
    """A plain paragraph lands in the placeholder as-is."""
    render(placeholder, "Bonjour")

    assert placeholder.text_frame.text == "Bonjour"


def test_empty_content_is_safe(placeholder):
    """An empty field clears the placeholder without failing."""
    render(placeholder, "")

    assert placeholder.text_frame.text == ""


def test_bold_and_italic_are_applied(placeholder):
    """Emphasis markers become run-level styles."""
    render(placeholder, "normal **gras** et *italique*")

    styles = {run.text: (run.font.bold, run.font.italic) for run in runs_of(placeholder)}
    assert styles["gras"] == (True, None)
    assert styles["italique"] == (None, True)


def test_unstyled_runs_inherit_from_the_template(placeholder):
    """Styles the Markdown does not ask for are left unset, not forced off."""
    render(placeholder, "**gras** puis normal")

    normal = next(run for run in runs_of(placeholder) if run.text.strip() == "puis normal")
    assert normal.font.bold is None
    assert normal.font.italic is None


def test_links_become_hyperlinks(placeholder):
    """A Markdown link carries its target onto the run."""
    render(placeholder, "voir [le site](https://example.org)")

    linked = next(run for run in runs_of(placeholder) if run.text == "le site")
    assert linked.hyperlink.address == "https://example.org"


def test_inline_code_uses_the_mono_font(placeholder):
    """Inline code switches font family and size."""
    render(placeholder, "appelle `ma_fonction`")

    code = next(run for run in runs_of(placeholder) if run.text == "ma_fonction")
    assert code.font.name == "Consolas"


def test_headings_are_rendered_bold(placeholder):
    """Headings have no pptx equivalent and degrade to a bold paragraph."""
    render(placeholder, "# Un titre")

    assert placeholder.text_frame.text == "Un titre"
    assert runs_of(placeholder)[0].font.bold is True


def test_nested_bullets_keep_their_level(placeholder):
    """List nesting maps onto paragraph levels."""
    render(placeholder, "- un\n- deux\n    - imbriqué")

    levels = [paragraph.level for paragraph in placeholder.text_frame.paragraphs]
    assert levels == [0, 0, 1]


def test_ordered_lists_are_numbered(placeholder):
    """Ordered items are prefixed with their number as plain text."""
    render(placeholder, "1. un\n2. deux")

    assert placeholder.text_frame.text == "1. un\n2. deux"


def test_nested_ordered_lists_restart_numbering(placeholder):
    """A nested ordered list numbers itself independently of its parent."""
    render(placeholder, "1. un\n2. deux\n    1. deux-a\n3. trois")

    assert placeholder.text_frame.text == "1. un\n2. deux\n1. deux-a\n3. trois"


def test_tables_are_rendered_as_tables(slide, placeholder):
    """A Markdown table becomes a pptx table with a bold header row."""
    render(placeholder, "| A | B |\n|---|---|\n| 1 | 2 |")

    table = table_of(slide)

    assert [[cell.text for cell in row.cells] for row in table.rows] == [["A", "B"], ["1", "2"]]
    assert table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold is True


def test_a_table_takes_over_the_whole_field(slide, placeholder):
    """Text alongside a table is dropped rather than rendered under it."""
    render(placeholder, "Du texte avant\n\n| A |\n|---|\n| 1 |")

    assert "Du texte avant" not in table_of(slide).cell(0, 0).text
