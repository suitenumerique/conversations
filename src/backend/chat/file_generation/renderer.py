"""Render Markdown into pptx shapes.

The model writes slide content as Markdown; this module turns it into formatted
runs, bulleted lists and tables inside a placeholder. Markdown is converted to
HTML first, then walked with BeautifulSoup, both of which the project already
depends on.

Inline styles are only ever set when the Markdown asks for them, never reset to
their falsy value, so the template's own styling keeps applying otherwise.
"""

from dataclasses import dataclass

import markdown as markdown_lib
from bs4 import BeautifulSoup, NavigableString, Tag
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.shapes.base import BaseShape
from pptx.text.text import TextFrame, _Paragraph
from pptx.util import Pt

BULLET_CHAR = "•"
MONO_FONT = "Consolas"
MONO_FONT_SIZE = Pt(10)

HEADING_TAGS = ("h1", "h2", "h3", "h4", "h5", "h6")
LIST_TAGS = ("ul", "ol")

# Left indent of a bullet level, and the hanging indent pulling the bullet
# character back into the margin, in EMU.
INDENT_PER_LEVEL = 342900
HANGING_INDENT = -342900


def render(target: BaseShape | TextFrame, markdown_text: str) -> None:
    """
    Render Markdown into a shape or a text frame.

    A Markdown table takes over the whole placeholder: when one is present, the
    rest of the content in the same field is dropped, because a pptx table is a
    separate graphic frame rather than something that can flow with text.
    """
    html = markdown_lib.markdown(markdown_text or "", extensions=["tables"])
    soup = BeautifulSoup(html, "html.parser")

    table = soup.find("table")
    if table is not None:
        # pylint: disable=protected-access
        shape = target if isinstance(target, BaseShape) else target._parent  # noqa: SLF001
        _render_table(shape, table)
        return

    text_frame = target.text_frame if hasattr(target, "text_frame") else target
    _render_text(text_frame, soup)


@dataclass
class _RenderContext:
    """Tracks the text frame being filled and whether its first paragraph is used."""

    text_frame: TextFrame
    first_paragraph_used: bool = False


@dataclass
class _InlineState:
    """Inline styles accumulated while walking down nested inline tags."""

    bold: bool = False
    italic: bool = False
    underline: bool = False
    mono: bool = False
    link: str | None = None


def _render_text(text_frame: TextFrame, soup: BeautifulSoup) -> None:
    """Fill a text frame with the block-level content of the parsed Markdown."""
    text_frame.clear()
    text_frame.word_wrap = True
    context = _RenderContext(text_frame=text_frame)

    for node in soup.contents:
        if isinstance(node, NavigableString) and not str(node).strip():
            continue

        name = getattr(node, "name", None)
        if name in HEADING_TAGS:
            paragraph = _new_paragraph(context)
            run = paragraph.add_run()
            run.text = node.get_text()
            run.font.bold = True
        elif name in LIST_TAGS:
            _render_list(context, node, level=0, ordered=name == "ol")
        else:
            paragraph = _new_paragraph(context)
            _render_inline_children(paragraph, node)


def _render_table(shape: BaseShape, table_node: Tag) -> None:
    """Replace a placeholder's content with a pptx table built from an HTML table."""
    rows = table_node.find_all("tr")
    if not rows:
        return

    row_count = len(rows)
    column_count = max(len(row.find_all(["th", "td"])) for row in rows)

    if hasattr(shape, "insert_table"):
        graphic_frame = shape.insert_table(rows=row_count, cols=column_count)
    else:
        graphic_frame = shape._parent.add_table(  # noqa: SLF001 # pylint: disable=protected-access
            row_count, column_count, shape.left, shape.top, shape.width, shape.height
        )

    table = graphic_frame.table
    for row_index, row_node in enumerate(rows):
        is_header_row = row_node.parent.name == "thead"
        for column_index, cell_node in enumerate(row_node.find_all(["th", "td"])):
            cell = table.cell(row_index, column_index)
            cell.text_frame.clear()
            cell.text_frame.word_wrap = True

            paragraph = cell.text_frame.paragraphs[0]
            paragraph.clear()
            _render_inline_children(paragraph, cell_node)

            if is_header_row or cell_node.name == "th":
                for run in paragraph.runs:
                    run.font.bold = True


def _new_paragraph(context: _RenderContext) -> _Paragraph:
    """Reuse the text frame's initial empty paragraph, then append new ones."""
    if not context.first_paragraph_used:
        paragraph = context.text_frame.paragraphs[0]
        paragraph.clear()
        context.first_paragraph_used = True
        return paragraph
    return context.text_frame.add_paragraph()


def _render_list(context: _RenderContext, list_node: Tag, level: int, ordered: bool) -> None:
    """Render a list and its nested lists, one paragraph per item."""
    items = list_node.find_all("li", recursive=False)
    for index, item in enumerate(items, start=1):
        _render_list_item(context, item, level, ordered, index)


def _render_list_item(
    context: _RenderContext, item: Tag, level: int, ordered: bool, index: int
) -> None:
    """Render one list item: its text into a marked paragraph, nested lists apart."""
    paragraph = None
    for child in item.contents:
        name = getattr(child, "name", None)

        if name in LIST_TAGS:
            _render_list(context, child, level + 1, ordered=name == "ol")
            continue

        if isinstance(child, NavigableString) and not str(child).strip():
            continue

        if paragraph is None:
            paragraph = _new_paragraph(context)
            _mark_list_paragraph(paragraph, level, ordered, index)

        _add_inline(paragraph, child)


def _mark_list_paragraph(paragraph: _Paragraph, level: int, ordered: bool, index: int) -> None:
    """Give a list item's paragraph its number or bullet."""
    if ordered:
        _apply_number(paragraph, level, index)
    else:
        _apply_bullet(paragraph, level)


def _render_inline_children(paragraph: _Paragraph, node: Tag) -> None:
    """Render every inline child of a block node into one paragraph."""
    for child in node.contents:
        _add_inline(paragraph, child)


def _add_inline(
    paragraph: _Paragraph,
    node: Tag | NavigableString,
    state: _InlineState | None = None,
) -> None:
    """Walk an inline subtree, emitting one run per text node with its styles."""
    state = state or _InlineState()

    if isinstance(node, NavigableString):
        if text := str(node):
            run = paragraph.add_run()
            run.text = text
            _apply_run_style(run, state)
        return

    tag = (node.name or "").lower()
    if tag == "br":
        return

    nested_state = _InlineState(
        bold=state.bold or tag in ("strong", "b"),
        italic=state.italic or tag in ("em", "i"),
        underline=state.underline or tag == "u",
        mono=state.mono or tag == "code",
        link=node.get("href") if tag == "a" else state.link,
    )
    for child in node.contents:
        _add_inline(paragraph, child, nested_state)


def _apply_run_style(run, state: _InlineState) -> None:
    """Apply the accumulated inline styles, leaving unset ones to the template."""
    if state.bold:
        run.font.bold = True
    if state.italic:
        run.font.italic = True
    if state.underline:
        run.font.underline = True
    if state.mono:
        run.font.name = MONO_FONT
        run.font.size = MONO_FONT_SIZE
    if state.link:
        run.hyperlink.address = state.link


def _paragraph_properties(paragraph: _Paragraph):
    """Return the paragraph's ``<a:pPr>`` element, creating it if needed."""
    # pylint: disable=protected-access
    return paragraph._p.get_or_add_pPr()  # noqa: SLF001


def _clear_bullet_properties(properties) -> None:
    """Drop any bullet definition inherited from the layout."""
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        element = properties.find(qn(tag))
        if element is not None:
            properties.remove(element)


def _apply_indent(properties, paragraph: _Paragraph, level: int) -> None:
    """Indent a paragraph to its nesting level, with a hanging first line."""
    properties.set("marL", str(INDENT_PER_LEVEL * (level + 1)))
    properties.set("indent", str(HANGING_INDENT))
    paragraph.level = level


def _apply_bullet(paragraph: _Paragraph, level: int) -> None:
    """Give an unordered list item its bullet character."""
    properties = _paragraph_properties(paragraph)
    _clear_bullet_properties(properties)
    _apply_indent(properties, paragraph, level)

    bullet = OxmlElement("a:buChar")
    bullet.set("char", BULLET_CHAR)
    properties.append(bullet)


def _apply_number(paragraph: _Paragraph, level: int, index: int) -> None:
    """
    Prefix an ordered list item with its number as plain text.

    Auto-numbering is not used on purpose: it restarts per shape and cannot be
    controlled per nesting level here, which produced wrong numbers on nested
    ordered lists.
    """
    properties = _paragraph_properties(paragraph)
    _clear_bullet_properties(properties)
    properties.append(OxmlElement("a:buNone"))
    _apply_indent(properties, paragraph, level)

    run = paragraph.add_run()
    run.text = f"{index}. "
