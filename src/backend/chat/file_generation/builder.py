"""Assemble a pptx deck from a validated presentation spec."""

import logging
from io import BytesIO

import pptx

from chat.file_generation.entities import Presentation, PresentationTemplate, Slide, SlideLayout
from chat.file_generation.renderer import render

logger = logging.getLogger(__name__)


class PresentationBuildError(Exception):
    """Raised when a deck cannot be assembled from its template."""


def build_presentation(template: PresentationTemplate, presentation: Presentation) -> bytes:
    """Render a presentation spec into the bytes of a pptx file."""
    try:
        deck = pptx.Presentation(str(template.path))
    except Exception as exc:
        logger.exception("Failed to open presentation template %s", template.path)
        raise PresentationBuildError(
            f"Could not open the presentation template at {template.path}."
        ) from exc

    for slide in presentation.slides:
        layout = template.get_layout(slide.type)
        if layout is None:
            # Unreachable for shipped types: the spec is validated against
            # SlideType and every member is bound (test_every_slide_type_...).
            # Guards a new type added without a binding.
            raise PresentationBuildError(f"No layout is bound to slide type {slide.type!r}.")
        _add_slide(deck, layout, slide)

    blob = BytesIO()
    deck.save(blob)
    return blob.getvalue()


def _add_slide(deck, layout: SlideLayout, slide: Slide) -> None:
    """Create a slide from its layout and fill the mapped placeholders."""
    pptx_layout = deck.slide_layouts.get_by_name(layout.layout_name)
    if pptx_layout is None:
        raise PresentationBuildError(
            f"Layout {layout.layout_name!r} is missing from the presentation template."
        )

    pptx_slide = deck.slides.add_slide(pptx_layout)
    for shape_index, field in layout.shape_index_to_field.items():
        render(pptx_slide.shapes[shape_index], getattr(slide, field, ""))

    if slide.slide_notes:
        render(pptx_slide.notes_slide.notes_text_frame, slide.slide_notes)


def check_template(template: PresentationTemplate) -> None:
    """
    Assert a template is consistent with the layouts declared against it.

    Called from the tests rather than at import time: it opens the file and
    builds a throwaway slide per layout, which is too costly to pay on every
    worker start.
    """
    deck = pptx.Presentation(str(template.path))
    layout_names = [layout.name for layout in deck.slide_layouts]

    for layout in template.layouts:
        pptx_layout = deck.slide_layouts.get_by_name(layout.layout_name)
        if pptx_layout is None:
            raise PresentationBuildError(
                f"Layout {layout.layout_name!r} not found in {template.path}. "
                f"Available layouts: {layout_names}"
            )

        # A layout and a slide created from it do not carry the same shapes, so
        # the declared indices can only be checked against a real slide.
        pptx_slide = deck.slides.add_slide(pptx_layout)
        shape_count = len(pptx_slide.shapes)
        if max(layout.shape_index_to_field) >= shape_count:
            raise PresentationBuildError(
                f"Layout {layout.layout_name!r} has {shape_count} shapes, but a "
                f"placeholder index {max(layout.shape_index_to_field)} is declared."
            )
